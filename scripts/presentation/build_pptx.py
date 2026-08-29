"""Builds the .pptx from the same rendered PNGs used for the PDF/previews
(guarantees visual match) plus real speaker notes per slide.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PNG_DIR = Path(sys.argv[1])
NOTES_PATH = Path(sys.argv[2])
OUT_PATH = Path(sys.argv[3])

notes = {n["stem"]: n for n in json.loads(NOTES_PATH.read_text())}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

png_files = sorted(PNG_DIR.glob("*.png"))
assert len(png_files) == 18, f"expected 18 PNGs, got {len(png_files)}"

for png in png_files:
    stem = png.stem
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)

    n = notes[stem]
    tf = slide.notes_slide.notes_text_frame
    tf.text = f"SAY: {n['say']}"
    tf.add_paragraph().text = f"EVIDENCE: {n['evidence']}"
    tf.add_paragraph().text = f"TRANSITION: {n['transition']}"
    tf.add_paragraph().text = f"LIKELY QUESTION: {n['question']}"
    tf.add_paragraph().text = f"ANSWER: {n['answer']}"

prs.save(str(OUT_PATH))
print(f"Wrote {OUT_PATH} ({len(png_files)} slides)")
